#!/usr/bin/env python3
"""Fill the CREATE-a-Thon presentation template with FoodBridge project content."""

from __future__ import annotations

from pathlib import Path
from typing import Optional

from pptx import Presentation

ROOT = Path(__file__).resolve().parent.parent
TEMPLATE = ROOT / "CREATE-a-Thon Agentic AI Solution.pptx"
OUTPUT = ROOT / "FoodBridge_CREATE-a-Thon_Presentation.pptx"


def replace_text(shape, new_text: str) -> None:
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    tf.paragraphs[0].text = new_text


def set_bullets(shape, lines: list[str], first_line: Optional[str] = None) -> None:
    tf = shape.text_frame
    tf.clear()
    all_lines = ([first_line] if first_line else []) + lines
    for i, line in enumerate(all_lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0


def build() -> Path:
    prs = Presentation(str(TEMPLATE))

    replace_text(prs.slides[0].shapes[0], "FoodBridge")
    replace_text(
        prs.slides[0].shapes[1],
        "Agentic AI for Food Waste Redistribution · CREATE-a-Thon",
    )

    replace_text(prs.slides[1].shapes[0], "Identification of Problem and Stakeholders")
    set_bullets(
        prs.slides[1].shapes[1],
        [
            "Problem: Restaurants, cafeterias, and grocery stores discard edible food "
            "while many people in Toronto face food insecurity.",
            "Gap: Municipal organic waste is high, yet surplus food rarely reaches "
            "shelters in time due to coordination friction.",
            "Stakeholders: Donors (restaurants/grocers), shelters & community kitchens, "
            "volunteer drivers, public health inspectors (Dinesafe), high-need communities.",
            "Data: Toronto Dinesafe health inspections + Canada Biomass municipal "
            "organic waste inventory (+ optional GDELT food-security signals).",
        ],
        first_line="Societal problem our agentic system addresses",
    )

    replace_text(prs.slides[2].shapes[0], "Ethical Considerations")
    set_bullets(
        prs.slides[2].shapes[1],
        [
            "Safety & liability → Pass-only Dinesafe donors; reject crucial (C) infraction history.",
            "Fairness in allocation → Small-organization recipients prioritized in matching queue.",
            "Transparency → Each agent logs decision steps (rule, input, outcome) in the dashboard.",
            "Privacy → Aggregate need data only; no individual beneficiary tracking.",
            "Small org exclusion → Quota-style priority prevents large hubs from dominating pickups.",
            "Environmental tradeoffs → Route optimizer vs. random routing disclosed; split runs if >25 km.",
            "Accountability → Ethics Guardian enforces human coordinator approval before dispatch.",
        ],
        first_line="Ethical issues and how FoodBridge addresses them",
    )

    replace_text(prs.slides[3].shapes[0], "Proof-of-Concept: Multi-Agent Architecture")
    set_bullets(
        prs.slides[3].shapes[1],
        [
            "Surplus Estimator — predicts daily kg surplus from Biomass + establishment type",
            "Need Prioritizer — ranks regions using GDELT + Biomass organic-waste signals",
            "Ethics & Donor Scout — screens Dinesafe Pass-only donors",
            "Matcher — pairs donors to nearest high-need shelters; fairness tiers",
            "Logistics Planner — nearest-neighbor pickup/delivery route",
            "Timing Negotiator — auto-negotiates evening pickup windows",
            "Ethics Guardian — final fairness audit & approval gate",
            "Deliverables: FastAPI web app (map + audit trail), CLI, REST API — github.com/kawsarahmedbhuiyan/food-bridge",
        ],
        first_line="Seven specialized agents, orchestrated with transparent decision logs",
    )

    replace_text(prs.slides[4].shapes[0], "Scalability & Impact")
    set_bullets(
        prs.slides[4].shapes[1],
        [
            "Geographic: Replace Dinesafe with any city's inspection open data; Biomass grids cover Canada.",
            "Technical: Modular agents — plug in live inventory APIs, OR-Tools routing, or SMS negotiation.",
            "Operational: Web dashboard + /api/plan for food-bank CRM and volunteer dispatch integration.",
            "Impact: Less edible waste to landfill, faster rescue windows, documented ethical decisions.",
            "Live demo: uvicorn web.app:app → select region → Run planning → map, matches, route, schedule.",
        ],
        first_line="How readily could this scale across industries and regions?",
    )

    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Created {path}")
